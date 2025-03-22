import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
from io import BytesIO
from typing import List, Dict, Tuple
import fitz  # PyMuPDF
from pptx.util import Inches
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import tempfile
from pptx2pdfwasm import PPTXtoPDFConverter
import sys
import asyncio


def detect_file_type(file):
    """Detecta si el archivo es PDF o PPTX"""
    file_type = file.name.split('.')[-1].lower()
    if file_type not in ['pdf', 'pptx']:
        st.error('Por favor, sube un archivo PDF o PPTX')
        return None
    return file_type

def get_file_stats(file_type, file):
    """Obtiene las estadísticas básicas del archivo"""
    try:
        if file_type == 'pdf':
            pdf_reader = PdfReader(file)
            return {
                "Tipo": "PDF",
                "Páginas": len(pdf_reader.pages),
            }
        else:  # pptx
            prs = Presentation(BytesIO(file.getvalue()))
            slides_with_notes = sum(1 for slide in prs.slides if slide.notes_slide and 
                                  slide.notes_slide.notes_text_frame.text.strip())
            return {
                "Tipo": "PowerPoint",
                "Diapositivas": len(prs.slides),
                "Con notas": slides_with_notes,
                "Sin notas": len(prs.slides) - slides_with_notes
            }
    except Exception as e:
        st.error(f"Error al procesar el archivo: {str(e)}")
        return None

def reset_state():
    """Reinicia el estado de la aplicación y elimina archivos temporales"""
    import os  # in case not already imported
    # Eliminar archivo de video generado
    generated_video = st.session_state.get("generated_video")
    if generated_video and os.path.exists(generated_video):
        try:
            os.remove(generated_video)
        except Exception as e:
            st.error(f"Error al eliminar el video temporal: {e}")
            
    keys_to_reset = [
        'step', 'uploaded_file', 'file_type', 'file_stats', 
        'generated_notes', 'target_language', 'tts_provider', 
        'generated_audio', 'video_options',
        'slides_images', 'slides_notes', 'user_prompt', 
        'vlm_model_url', 'vlm_model_id', 'slides_notes_audios', 
        'selected_voice', 'generated_video'
    ]
    for key in keys_to_reset:
        if key in st.session_state:
            del st.session_state[key]

def init_session_state():
    """Inicializa las variables de estado necesarias"""
    defaults = {
        'step': 0,
        'uploaded_file': None,
        'file_type': None,
        'file_stats': None,
        'generated_notes': None,
        'target_language': 'es',
        'tts_provider': 'elevenlabs',
        'generated_audio': None,
        'video_options': {'fps': 1, 'transition': 'fade'}
    }
    
    for key, default_value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = default_value

def get_language_options() -> Dict[str, str]:
    """Retorna un diccionario de idiomas disponibles y codigo BCP-47"""
    return {
        'es': 'Español 🇪🇸',
        'en': 'English 🇬🇧',
        'fr': 'Français 🇫🇷',
        'de': 'Deutsch 🇩🇪',
        'it': 'Italiano 🇮🇹',
        'pt': 'Português 🇵🇹',
        'ar': 'العربية 🇸🇦',
    }

def extract_slides_content(presentation: Presentation) -> List[dict]:
    """Extrae el contenido y notas de las diapositivas"""
    slides_data = []
    for idx, slide in enumerate(presentation.slides, 1):
        # Extraer texto de la diapositiva
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        # Extraer notas
        notes = ""
        if slide.notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            notes = slide.notes_slide.notes_text_frame.text.strip()
        
        slides_data.append({
            'slide_num': idx,
            'content': ' '.join(slide_text),
            'notes': notes
        })
    
    return slides_data

def extract_pdf_slides(file) -> List[bytes]:
    """Extrae las diapositivas de un archivo PDF como imágenes usando PyMuPDF"""
    images = []
    try:
        file.seek(0)  # reset the file pointer to the beginning
        pdf_document = fitz.open(stream=file.read(), filetype="pdf")
        for page_number in range(pdf_document.page_count):
            page = pdf_document.load_page(page_number)
            pix = page.get_pixmap()
            img_data = pix.tobytes("png")  # Convert to PNG format
            images.append(img_data)
    except Exception as e:
        st.error(f"Error al procesar el PDF con PyMuPDF: {e}")
        return []
    return images

def extract_pptx_slides(uploaded_file):
    """
    Extrae imágenes y notas de cada slide de un archivo PPTX.
    
    Las imágenes se obtienen convirtiendo el PPTX a PDF con pptx2pdfwasm
    y extrayendo cada página del PDF en formato PNG.
    Las notas se extraen utilizando python-pptx.
    
    Args:
        uploaded_file (BytesIO): Archivo PPTX subido.
    
    Returns:
        Tuple[List[bytes], List[str]]: (lista de imágenes en bytes, lista de notas)
    """
    slides_images = []
    slides_notes = []
    pptx_path = None
    pdf_path = None

    try:
        # Guardar el archivo PPTX en un archivo temporal
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(uploaded_file.read())
            pptx_path = tmp.name
        
        # Asegurar que el archivo temporal se cierre antes de la conversión

        if sys.platform.startswith('win') and sys.version_info >= (3, 8):
            asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())
        
        # Convertir el PPTX a PDF usando pptx2pdfwasm
        converter = PPTXtoPDFConverter(headless=True, log_enabled=False, port=7777)
        converter.start_server()
        try:
            converter.convert(pptx_path, os.path.splitext(pptx_path)[0] + ".pdf")
        finally:
            converter.stop_server()

        # Definir la ruta del PDF generado
        pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
        if not os.path.exists(pdf_path):
            st.error(f"No se encontró el PDF generado con pptx2pdfwasm en: {pdf_path}")
            return [], []

        # Abrir el PDF y extraer cada página como imagen PNG
        doc = fitz.open(pdf_path)
        for page in doc:
            pix = page.get_pixmap()
            slides_images.append(pix.tobytes("png"))
        doc.close()

        # Extraer las notas de cada slide usando python-pptx
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            slides_notes.append(notes)

    except Exception as e:
        st.error(f"Error al procesar el archivo PPTX: {e}")
        return [], []
    finally:
        # Limpiar archivos temporales
        try:
            if pptx_path and os.path.exists(pptx_path):
                os.unlink(pptx_path)
            if pdf_path and os.path.exists(pdf_path):
                os.unlink(pdf_path)
        except Exception as e:
            st.error(f"Error al limpiar archivos temporales: {e}")

    return slides_images, slides_notes

def get_file_bytes(file_path: str) -> bytes:
    with open(file_path, "rb") as f:
        return f.read()
